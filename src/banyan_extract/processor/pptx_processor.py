import re
import io
import PIL
from typing import List, Tuple, Union

from tqdm import tqdm

try:
    from surya.texify import TexifyPredictor # May need to find an alternative
    USE_OCR = True
except:
    USE_OCR = False
    print("Surya ocr not installed; ocr not available for pptx processing")

from pptx import Presentation
from pptx.shapes.group import GroupShape

from .processor import Processor
from ..output.pptx_output import PptxOutput
from ..ocr.nemotron_ocr import NemotronOCR


try:
    class MarkdownTexifyPredictor(TexifyPredictor):
        def fix_fences(self, text: str) -> str:
            text = re.sub(r'<math display="block">(.*?)</math>',r'$$\1$$', text, flags=re.DOTALL)
            text = re.sub(r'<math>(.*?)</math>',r'$\1$', text, flags=re.DOTALL)
            if re.search(r'<math display="block">', text):
                text = ""
            if re.search(r'<math>', text):
                text = ""
            return text
except:
    class MarkdownTexifyPredictor:
        pass


class PptxProcessor(Processor):
    
    def __init__(self, ocr_backend="surya", nemotron_endpoint="", nemotron_model="nvidia/nemoretriever-parse"):
        """
        Initialize PPTX processor with OCR backend selection.
        
        Args:
            ocr_backend: Which OCR backend to use ('surya' or 'nemotron')
            nemotron_endpoint: URL for Nemotron parse endpoint (if using nemotron)
            nemotron_model: Model name for Nemotron OCR
        """
        super().__init__()
        self.ocr_backend = None  # Single member variable for OCR backend
        self.ocr_available = False

        # Initialize selected OCR backend
        if ocr_backend == "nemotron":
            try:
                self.ocr_backend = NemotronOCR(
                    endpoint_url=nemotron_endpoint,
                    model_name=nemotron_model
                )
                self.ocr_available = True
            except Exception as e:
                print(f"Failed to initialize Nemotron OCR: {e}")
                self.ocr_available = False
        else:  # Default to Surya
            try:
                self.ocr_backend = MarkdownTexifyPredictor()
                self.ocr_available = USE_OCR
            except:
                self.ocr_available = False
                print("Surya OCR not installed; OCR not available for PPTX processing")

    def ocr_image(self, image):
        """Perform OCR using the selected backend."""
        if not self.ocr_available or self.ocr_backend is None:
            return ""

        try:
            if isinstance(self.ocr_backend, NemotronOCR):
                # Use Nemotron OCR
                return self.ocr_backend.ocr_image(image)
            else:
                # Use Surya OCR
                ocr_output = self.ocr_backend([image])[0]
                if ocr_output.text is not None:
                    return ocr_output.text
        except Exception as e:
            print(f"OCR failed: {e}")

        return ""

    def process_image(self, image):
        if 'wmf' in image.content_type:
            print("pptx contains a wmf image, skipping image")
            return None
        else:
            byte_io = io.BytesIO(image.blob)
            image_object = PIL.Image.open(byte_io)
            return image_object

    def process_document(self, filepath, rotation_angle: Union[int, float] = 0):
        # Note: PPTX processor doesn't currently support rotation as it works with structured slides
        # For future implementation, we would need to rotate individual images within slides
        if rotation_angle != 0:
            print(f"Warning: Rotation is not currently supported for PptxProcessor. Angle {rotation_angle} will be ignored.")
        
        prs = Presentation(filepath)
        images = []
        slide_texts = []
        for slide in prs.slides:
            slide_text = []
            images.append([])
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for par in shape.text_frame.paragraphs:
                        slide_text.append(par.text)
                if "PIC" in str(shape.shape_type):
                    image = self.process_image(shape.image)
                    if image is not None:
                        image_ocr = self.ocr_image(image)
                        slide_text.append(image_ocr)
                        images[-1].append(image)
                if "GROUP" in str(shape.shape_type):
                    for sub_shape in shape.shapes:
                        if "PIC" in str(sub_shape.shape_type):
                            image = self.process_image(sub_shape.image)
                            if image is not None:
                                image_ocr = self.ocr_image(image)
                                slide_text.append(image_ocr)
                                images[-1].append(image)

                        if sub_shape.has_text_frame:
                            for par in sub_shape.text_frame.paragraphs:
                                slide_text.append(par.text)
            slide_texts.append("\n".join(slide_text))

        metadata = {}
        return PptxOutput(slide_texts, images, metadata)

    def process_batch_documents(self, filepaths, rotation_angle: Union[int, float] = 0):
        # Note: Batch processing not currently implemented for PPTX processor
        if rotation_angle != 0:
            print(f"Warning: Rotation is not currently supported for PptxProcessor. Angle {rotation_angle} will be ignored.")
        
        # For future implementation, we would process each file with rotation support
        pass
